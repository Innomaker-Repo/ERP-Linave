from pptx import Presentation

prs = Presentation('Rascunho_gabriel.pptx')

# Delete existing slides so only the new SCA ones are in the final file
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
for slide in slides:
    xml_slides.remove(slide)

slides_data = [
    {
        "layout": "Title 1 Dark",
        "title": "Software Composition Analysis (SCA)",
        "content": "Gestão de Dependências, Supply Chain Security e Remediação Automatizada no Checkmarx One",
        "placeholders": {10: "title", 11: "content"}
    },
    {
        "layout": "1_Title & Content",
        "title": "Sessão 5 - SCA e O Desafio Atual",
        "content": "• O desenvolvimento moderno depende enormemente do uso de pacotes Open Source (cerca de 70-90% do código).\n• Vetor de Ataque: Bibliotecas de terceiros herdam os privilégios da aplicação. Uma falha expõe o ambiente.\n• Exploitable Path: O Checkmarx correlaciona SAST e SCA para validar se o código-fonte de fato invoca a função vulnerável, reduzindo falsos positivos.",
        "placeholders": {0: "title", 1: "content"}
    },
    {
        "layout": "1_Two Columns",
        "title": "Análise de Dependências",
        "content1": "Dependências Diretas",
        "subcontent1": "Bibliotecas declaradas explicitamente pela equipe no projeto (ex: package.json).",
        "content2": "Dependências Transitivas",
        "subcontent2": "Pacotes exigidos internamente pelas dependências diretas. 80% das vulnerabilidades residem aqui (o ponto cego). O Checkmarx mapeia toda a topologia.",
        "placeholders": {0: "title", 16: "content1", 17: "subcontent1", 18: "content2", 19: "subcontent2"}
    },
    {
        "layout": "1_Title & Content",
        "title": "Compliance e Riscos Legais",
        "content": "• O uso não governado de software livre gera altos riscos corporativos.\n• Licenças Permissivas (MIT, Apache 2.0): Baixo risco, uso amplo com poucas restrições.\n• Licenças Copyleft (GPL): Alto risco para código fechado. Exige que derivados também sejam open source.\n• O tenant do Checkmarx alerta ativamente no painel de 'Legal Risk' sobre infrações de políticas.",
        "placeholders": {0: "title", 1: "content"}
    },
    {
        "layout": "1_Three Columns",
        "title": "Supply Chain Security: Malicious Packages",
        "content1": "Reputação",
        "subcontent1": "Análise do mantenedor. Prevê riscos como 'New User' e 'Account Takeover / Repojacking' (roubo de credenciais).",
        "content2": "Confiabilidade",
        "subcontent2": "Anomalias de hospedagem. Prevê 'Typosquatting' (nomes parecidos, ex: react-doms) e 'Dependency Confusion'.",
        "content3": "Comportamento",
        "subcontent3": "Análise da função. Evita Data Exfiltration (DNS Tunneling) e injeção de Crypto Miners.",
        "placeholders": {0: "title", 10: "content1", 13: "subcontent1", 11: "content2", 14: "subcontent2", 12: "content3", 15: "subcontent3"}
    },
    {
        "layout": "1_Two Columns",
        "title": "Eficiência e Visibilidade Global",
        "content1": "Software Bill of Materials (SBOM)",
        "subcontent1": "Exportação nativa do 'recibo de ingredientes' nos padrões globais CycloneDX e SPDX para compliance. Pesquisa cruzada em portfólio em casos Zero-Day.",
        "content2": "SCA Scan Recalculation",
        "subcontent2": "Resultados atualizados sem novos scans de pipeline se: novas ameaças forem descobertas no BD, ações de triage forem feitas ou políticas mudarem.",
        "placeholders": {0: "title", 16: "content1", 17: "subcontent1", 18: "content2", 19: "subcontent2"}
    },
    {
        "layout": "1_Title & Content",
        "title": "Auditoria Inteligente: Padrão VEX",
        "content": "• O Vulnerability Exploitability eXchange (VEX) complementa o SBOM.\n• Problema: Um pacote pode estar vulnerável, mas a falha não ser explorável no seu contexto, gerando ruído de alerta.\n• Solução: O VEX anexa o 'State' (ex: Not Affected) e a 'Justification' (ex: Code not reachable).\n• O Checkmarx exporta o status VEX embutido automatizadamente no CycloneDX.",
        "placeholders": {0: "title", 1: "content"}
    },
    {
         "layout": "1_Two Columns",
        "title": "Automação e Qualidade em CI/CD",
        "content1": "Remediação via Auto PR",
        "subcontent1": "O AppSec Knowledge Center sugere o version bump exato. A ferramenta gera e envia automaticamente o Pull Request de correção (suporte a GitHub, GitLab, Azure).",
        "content2": "Break Build / Quality Gates",
        "subcontent2": "Como você está orquestrando os scans no GitLab CI/CD, o foco é o Gate de Qualidade. Se limites de política forem violados, ele encerra a esteira com falha ativa.",
        "placeholders": {0: "title", 16: "content1", 17: "subcontent1", 18: "content2", 19: "subcontent2"}
    }
]

layout_map = {layout.name: layout for layout in prs.slide_layouts}

for slide_def in slides_data:
    layout_name = slide_def["layout"]
    if layout_name in layout_map:
        slide_layout = layout_map[layout_name]
        slide = prs.slides.add_slide(slide_layout)
        for ph_idx, content_key in slide_def["placeholders"].items():
            try:
                shape = slide.placeholders[ph_idx]
                shape.text = slide_def[content_key]
            except KeyError:
                 pass

prs.save('SCA_Checkmarx_Apresentacao.pptx')
print("File regenerated.")